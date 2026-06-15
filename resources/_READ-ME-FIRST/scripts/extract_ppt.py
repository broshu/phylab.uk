#!/usr/bin/env python3
"""Extract images and a text outline from a .pptx.

Usage:
    python3 extract_ppt.py <source.pptx> <out_img_dir>

Writes every picture to <out_img_dir>/slideN_picM.<ext>, a manifest.json
describing them, and prints the slide-by-slide text outline to stdout so you can
reconstruct the teaching logic. Requires python-pptx (pip install python-pptx).
"""
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx --break-system-packages")


def walk(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(sh.shapes)


def main():
    if len(sys.argv) != 3:
        sys.exit(__doc__)
    src, out = Path(sys.argv[1]), Path(sys.argv[2])
    out.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(src))

    manifest, seen = [], {}
    print(f"# Outline: {src.name}  ({len(prs.slides)} slides)\n")
    for i, slide in enumerate(prs.slides, 1):
        texts, pic_n = [], 0
        for sh in walk(slide.shapes):
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    texts.append(t)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_n += 1
                img = sh.image
                key = hash(img.blob)
                ext = img.ext or "png"
                name = f"slide{i}_pic{pic_n}.{ext}"
                (out / name).write_bytes(img.blob)
                entry = {
                    "file": name, "slide": i,
                    "w_emu": getattr(sh, "width", None),
                    "h_emu": getattr(sh, "height", None),
                    "duplicate_of": seen.get(key),
                }
                seen.setdefault(key, name)
                manifest.append(entry)
        print(f"## Slide {i}")
        for t in texts:
            for line in t.splitlines():
                if line.strip():
                    print(f"   - {line.strip()}")
        if pic_n:
            print(f"   [{pic_n} image(s)]")
        print()

    (out / "manifest.json").write_text(json.dumps(manifest, indent=2, ensure_ascii=False))
    print(f"\nExtracted {len(manifest)} image(s) -> {out}/  (see manifest.json)")
    dups = [m for m in manifest if m["duplicate_of"]]
    if dups:
        print(f"Note: {len(dups)} are duplicates of earlier images.")


if __name__ == "__main__":
    main()
